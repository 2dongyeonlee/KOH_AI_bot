import os
import io
import fitz  # PyMuPDF
from pptx import Presentation
from google import genai
from google.genai import types

SYSTEM_PROMPT = """너는 SK하이닉스 커뮤니케이션 총괄 조직의
6R 전략 담당 권오혁 담당의 전담 비서 '궁오혁봇'이야.

권오혁 담당은 국내외 대외 이슈를 즉각 포착하고,
입체적인 커뮤니케이션 메시지를 적재적소에 배치하며,
높은 정무감각으로 SK하이닉스의 커뮤니케이션을
전략적으로 대응하는 전문가야.

행동 원칙:
- 답변은 항상 간결하고 실용적으로
- 존댓말 유지
- 커뮤니케이션·IR·GR·PR·ER·CR·BR 관련 질문에 전문적으로 답변
- 불필요한 부연 설명 없이 바로 본론으로"""

SUMMARY_PROMPT = "다음 내용을 핵심만 3줄로 요약해줘.\n각 줄은 번호를 붙여줘. 한국어로.\n\n{content}"

EXT_TO_MIME = {
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "png": "image/png",
    "gif": "image/gif",
    "webp": "image/webp",
    "bmp": "image/bmp",
}


def _client() -> genai.Client:
    return genai.Client(api_key=os.environ["GEMINI_API_KEY"])


def _config() -> types.GenerateContentConfig:
    return types.GenerateContentConfig(system_instruction=SYSTEM_PROMPT)


def summarize_text(text: str) -> str:
    try:
        response = _client().models.generate_content(
            model="gemini-2.0-flash",
            contents=SUMMARY_PROMPT.format(content=text),
            config=_config(),
        )
        return response.text
    except Exception:
        return "요약 중 오류가 발생했어요. 잠시 후 다시 시도해주세요."


def summarize_pdf(file_bytes: bytes) -> str:
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = "\n".join(page.get_text() for page in doc)
        if not text.strip():
            return "문서를 읽을 수 없어요. 텍스트가 포함된 PDF인지 확인해주세요."
        return summarize_text(text)
    except Exception:
        return "문서를 읽을 수 없어요. 텍스트가 포함된 PDF인지 확인해주세요."


def summarize_image(file_bytes: bytes, ext: str = "jpg") -> str:
    try:
        mime = EXT_TO_MIME.get(ext, "image/jpeg")
        image_part = types.Part.from_bytes(data=file_bytes, mime_type=mime)
        prompt = "이 이미지의 내용을 핵심만 3줄로 요약해줘.\n각 줄은 번호를 붙여줘. 한국어로."
        response = _client().models.generate_content(
            model="gemini-2.0-flash",
            contents=[image_part, prompt],
            config=_config(),
        )
        return response.text
    except Exception:
        return "요약 중 오류가 발생했어요. 잠시 후 다시 시도해주세요."


def summarize_pptx(file_bytes: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
        text = "\n".join(texts)
        if not text.strip():
            return "문서를 읽을 수 없어요. 텍스트가 포함된 PPTX인지 확인해주세요."
        return summarize_text(text)
    except Exception:
        return "문서를 읽을 수 없어요. 텍스트가 포함된 PPTX인지 확인해주세요."
